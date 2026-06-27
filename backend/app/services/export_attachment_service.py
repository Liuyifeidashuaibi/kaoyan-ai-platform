"""Build export attachments (txt / docx / pdf) with commercial-grade formatting.

PDF: uses reportlab SimpleDocTemplate + Paragraph with CID font STSong-Light (body)
     and TTF SimHei (headings/bold). Includes cover page, page numbers, proper margins.
DOCX: uses python-docx with SimSun (body) and SimHei (headings), cover page,
      page numbers, first-line indent, proper margins.
"""

from __future__ import annotations

import logging
import re
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any

from app.config import get_settings

logger = logging.getLogger(__name__)

MIME_BY_FORMAT = {
    "txt": "text/plain",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


# ── Utilities ──────────────────────────────────────────────


def _safe_stem(title: str) -> str:
    stem = "".join(ch if ch.isalnum() or ch in " -_" else "_" for ch in title.strip())
    return (stem.strip(" _-") or "document")[:80]


def _escape_xml(text: str) -> str:
    """Escape XML special characters for reportlab Paragraph."""
    text = text.replace("&", "&amp;")
    text = text.replace("<", "&lt;")
    text = text.replace(">", "&gt;")
    return text


def _convert_inline_for_pdf(text: str) -> str:
    """Convert **bold**, `code` to reportlab Paragraph XML tags."""
    text = _escape_xml(text)
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
    text = re.sub(r"`(.+?)`", r'<font face="STSong-Light">\1</font>', text)
    return text


# ── Markdown line parser ───────────────────────────────────


def _parse_markdown_line(line: str) -> tuple[str, str, int]:
    """
    Parse a single Markdown line.
    Returns (line_type, text, level).
      line_type: "h1"|"h2"|"h3"|"bullet"|"numbered"|"code_fence"|"quote"|"text"
    """
    stripped = line.strip()

    if stripped.startswith("```"):
        return ("code_fence", stripped[3:].strip(), 0)
    if stripped.startswith("### "):
        return ("h3", stripped[4:].strip(), 3)
    if stripped.startswith("## "):
        return ("h2", stripped[3:].strip(), 2)
    if stripped.startswith("# "):
        return ("h1", stripped[2:].strip(), 1)
    if stripped.startswith("- ") or stripped.startswith("* "):
        return ("bullet", stripped[2:].strip(), 0)
    if stripped.startswith("  - ") or stripped.startswith("  * "):
        return ("bullet", stripped[4:].strip(), 1)
    m = re.match(r"^(\d+)\.\s+(.+)", stripped)
    if m:
        return ("numbered", m.group(2).strip(), int(m.group(1)))
    if stripped.startswith("> "):
        return ("quote", stripped[2:].strip(), 0)
    return ("text", stripped, 0)


# ── Font registration ──────────────────────────────────────


def _register_pdf_fonts() -> tuple[str, str]:
    """
    Register fonts for PDF export.
    Returns (body_font, heading_font).

    Body: STSong-Light (CID, always works with Chinese)
    Heading: SimHei (TTF from system, for bold/headings)
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.pdfmetrics import registerFontFamily

    body = "STSong-Light"
    heading = body  # fallback: same as body

    # Register CID body font (always available in reportlab)
    if body not in pdfmetrics.getRegisteredFontNames():
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(body))
        except Exception as e:
            logger.warning("Failed to register CID font %s: %s", body, e)

    # Try TTF heading font from system
    settings = get_settings()
    ttf_candidates = [
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSansSC-Bold.ttf"),
        settings.root / "data" / "fonts" / "NotoSansSC-Regular.ttf",
    ]
    for path in ttf_candidates:
        if path.is_file():
            try:
                name = "SimHei" if "simhei" in path.name.lower() else "NotoSansSC"
                if name not in pdfmetrics.getRegisteredFontNames():
                    from reportlab.pdfbase.ttfonts import TTFont
                    pdfmetrics.registerFont(TTFont(name, str(path)))
                heading = name
                logger.info("PDF heading font: %s from %s", name, path)
                break
            except Exception as e:
                logger.warning("Failed to register TTF font %s: %s", path, e)
                continue

    # Register font family so <b> tags use heading font
    try:
        registerFontFamily(
            body, normal=body, bold=heading, italic=body, boldItalic=heading
        )
    except Exception:
        pass

    return body, heading


# ── TXT ────────────────────────────────────────────────────


def build_txt(content: str, title: str, author_info: str = "") -> tuple[bytes, str, str]:
    body = content.replace("\r\n", "\n")
    header = title.strip()
    if author_info and author_info.strip():
        header = f"{header}\n{author_info.strip()}"
    if header:
        body = f"{header}\n\n{body}"
    data = body.encode("utf-8")
    return data, f"{_safe_stem(title)}.txt", MIME_BY_FORMAT["txt"]


# ── PDF ────────────────────────────────────────────────────


def build_pdf(content: str, title: str, author_info: str = "") -> tuple[bytes, str, str]:
    """Build a commercial-grade PDF with cover page, page numbers, and proper formatting."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        PageBreak,
    )

    body_font, heading_font = _register_pdf_fonts()

    # ── Paragraph styles ──
    styles = {
        "cover_title": ParagraphStyle(
            "CoverTitle", fontName=heading_font, fontSize=28,
            alignment=TA_CENTER, leading=42, spaceAfter=20,
        ),
        "cover_info": ParagraphStyle(
            "CoverInfo", fontName=body_font, fontSize=14,
            alignment=TA_CENTER, leading=21, spaceAfter=8,
        ),
        "h1": ParagraphStyle(
            "H1", fontName=heading_font, fontSize=18,
            spaceBefore=16, spaceAfter=10, leading=26,
        ),
        "h2": ParagraphStyle(
            "H2", fontName=heading_font, fontSize=15,
            spaceBefore=14, spaceAfter=8, leading=22,
        ),
        "h3": ParagraphStyle(
            "H3", fontName=heading_font, fontSize=13,
            spaceBefore=12, spaceAfter=6, leading=19,
        ),
        "body": ParagraphStyle(
            "Body", fontName=body_font, fontSize=11,
            alignment=TA_JUSTIFY, firstLineIndent=22,
            spaceBefore=3, spaceAfter=3, leading=18,
        ),
        "bullet": ParagraphStyle(
            "Bullet", fontName=body_font, fontSize=11,
            leftIndent=20, spaceBefore=2, spaceAfter=2, leading=18,
        ),
        "code": ParagraphStyle(
            "Code", fontName=body_font, fontSize=9,
            leftIndent=16, leading=14,
            backColor="#f5f5f5", borderPadding=6,
        ),
        "quote": ParagraphStyle(
            "Quote", fontName=body_font, fontSize=11,
            leftIndent=24, textColor="#555555",
            spaceBefore=4, spaceAfter=4, leading=18,
        ),
    }

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2.5 * cm, rightMargin=2.5 * cm,
        topMargin=2.5 * cm, bottomMargin=2.5 * cm,
        title=title, author="考研AI平台",
    )

    story: list = []

    # ── Cover page ──
    story.append(Spacer(1, 5 * cm))
    story.append(Paragraph(_escape_xml(title.strip()), styles["cover_title"]))
    story.append(Spacer(1, 2 * cm))
    if author_info and author_info.strip():
        for line in author_info.strip().split("\n"):
            if line.strip():
                story.append(Paragraph(_escape_xml(line.strip()), styles["cover_info"]))
    else:
        story.append(Paragraph(
            datetime.now().strftime("%Y年%m月%d日"), styles["cover_info"]
        ))
    story.append(PageBreak())

    # ── Parse content ──
    lines = content.replace("\r\n", "\n").split("\n")
    in_code = False
    code_buf: list[str] = []

    for line in lines:
        ltype, ltext, level = _parse_markdown_line(line)

        if ltype == "code_fence":
            if in_code:
                code_text = _escape_xml("\n".join(code_buf)).replace("\n", "<br/>")
                story.append(Paragraph(code_text, styles["code"]))
                story.append(Spacer(1, 6))
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_buf.append(line)
            continue

        if not ltext and ltype == "text":
            story.append(Spacer(1, 6))
            continue

        if ltype == "h1":
            story.append(Paragraph(_convert_inline_for_pdf(ltext), styles["h1"]))
        elif ltype == "h2":
            story.append(Paragraph(_convert_inline_for_pdf(ltext), styles["h2"]))
        elif ltype == "h3":
            story.append(Paragraph(_convert_inline_for_pdf(ltext), styles["h3"]))
        elif ltype == "bullet":
            indent = "&nbsp;&nbsp;" if level > 0 else ""
            story.append(Paragraph(
                f"{indent}• {_convert_inline_for_pdf(ltext)}", styles["bullet"]
            ))
        elif ltype == "numbered":
            story.append(Paragraph(
                f"{level}. {_convert_inline_for_pdf(ltext)}", styles["bullet"]
            ))
        elif ltype == "quote":
            story.append(Paragraph(_convert_inline_for_pdf(ltext), styles["quote"]))
        else:
            story.append(Paragraph(_convert_inline_for_pdf(ltext), styles["body"]))

    # Close unclosed code block
    if in_code and code_buf:
        code_text = _escape_xml("\n".join(code_buf)).replace("\n", "<br/>")
        story.append(Paragraph(code_text, styles["code"]))

    # ── Page number callback ──
    def _on_page(canvas, doc):
        canvas.saveState()
        canvas.setFont(body_font, 9)
        page_num = canvas.getPageNumber()
        canvas.drawCentredString(A4[0] / 2, 1.5 * cm, f"— {page_num} —")
        canvas.restoreState()

    doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)

    return buf.getvalue(), f"{_safe_stem(title)}.pdf", MIME_BY_FORMAT["pdf"]


# ── DOCX ───────────────────────────────────────────────────


def _set_run_font(run, font_name: str, size=None, bold: bool = False):
    """Set font name and East Asian font for a docx run."""
    run.font.name = font_name
    if size:
        run.font.size = size
    if bold:
        run.bold = True
    from docx.oxml.ns import qn
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn("w:eastAsia"), font_name)


def _add_docx_runs(paragraph, text: str):
    """Add text to docx paragraph with **bold** and `code` inline formatting."""
    parts = re.split(r"(\*\*.+?\*\*|`.+?`)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            _set_run_font(paragraph.add_run(part[2:-2]), "SimSun", bold=True)
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Consolas"
        else:
            _set_run_font(paragraph.add_run(part), "SimSun")


def build_docx(content: str, title: str, author_info: str = "") -> tuple[bytes, str, str]:
    """Build a commercial-grade DOCX with cover page, page numbers, and proper formatting."""
    from docx import Document
    from docx.shared import Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Page margins ──
    for section in doc.sections:
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(3.18)
        section.right_margin = Cm(3.18)

    # ── Default font: SimSun ──
    style = doc.styles["Normal"]
    style.font.name = "SimSun"
    style.font.size = Pt(12)
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn("w:eastAsia"), "SimSun")

    # ── Cover page ──
    for _ in range(6):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _set_run_font(p.add_run(title.strip()), "SimHei", Pt(28), bold=True)

    if author_info and author_info.strip():
        for line in author_info.strip().split("\n"):
            if line.strip():
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                _set_run_font(p.add_run(line.strip()), "SimSun", Pt(14))
    else:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _set_run_font(
            p.add_run(datetime.now().strftime("%Y年%m月%d日")), "SimSun", Pt(14)
        )

    doc.add_page_break()

    # ── Parse content ──
    lines = content.replace("\r\n", "\n").split("\n")
    in_code = False
    code_buf: list[str] = []

    for line in lines:
        ltype, ltext, level = _parse_markdown_line(line)

        if ltype == "code_fence":
            if in_code:
                p = doc.add_paragraph()
                _set_run_font(p.add_run("\n".join(code_buf)), "SimSun", Pt(10))
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_buf.append(line)
            continue

        if not ltext and ltype == "text":
            doc.add_paragraph()
            continue

        if ltype == "h1":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(16)
            p.paragraph_format.space_after = Pt(10)
            _set_run_font(p.add_run(ltext), "SimHei", Pt(18), bold=True)
        elif ltype == "h2":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(14)
            p.paragraph_format.space_after = Pt(8)
            _set_run_font(p.add_run(ltext), "SimHei", Pt(15), bold=True)
        elif ltype == "h3":
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(6)
            _set_run_font(p.add_run(ltext), "SimHei", Pt(13), bold=True)
        elif ltype == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            _add_docx_runs(p, ltext)
        elif ltype == "numbered":
            p = doc.add_paragraph(style="List Number")
            _add_docx_runs(p, ltext)
        elif ltype == "quote":
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(1)
            _add_docx_runs(p, ltext)
        else:
            p = doc.add_paragraph()
            p.paragraph_format.first_line_indent = Cm(0.74)
            _add_docx_runs(p, ltext)

    # Close unclosed code block
    if in_code and code_buf:
        p = doc.add_paragraph()
        _set_run_font(p.add_run("\n".join(code_buf)), "SimSun", Pt(10))

    # ── Page numbers in footer ──
    section = doc.sections[0]
    footer = section.footer
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("— ")
    # PAGE field
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.text = "PAGE"
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run._element.append(fld_begin)
    run._element.append(instr)
    run._element.append(fld_end)
    p.add_run(" —")

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue(), f"{_safe_stem(title)}.docx", MIME_BY_FORMAT["docx"]


# ── Excel ───────────────────────────────────────────────


def build_excel(
    sheets_data: list[dict],
    title: str,
) -> tuple[bytes, str]:
    """
    Build an Excel (.xlsx) file with formatted sheets.

    Args:
        sheets_data: List of sheet definitions.
            Each sheet: {"name": str, "headers": list[str], "rows": list[list]}
    Returns:
        (file_bytes, filename)
    """
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    # Remove default sheet if we have custom ones
    if len(wb.worksheets) > 0 and sheets_data:
        wb.remove(wb.active)

    # Styles
    header_font = Font(name="Microsoft YaHei", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    body_font = Font(name="Microsoft YaHei", size=10)
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)

    for sheet_def in sheets_data:
        sheet_name = sheet_def.get("name", f"Sheet{len(wb.worksheets) + 1}")[:31]  # Excel 31 char limit
        ws = wb.create_sheet(title=sheet_name)

        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])

        # Write headers
        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=str(header))
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = center_align
            cell.border = thin_border

        # Write data rows
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                # Try to convert numeric strings to numbers
                cell_value: Any = value
                if isinstance(value, str):
                    try:
                        if "." in value:
                            cell_value = float(value)
                        else:
                            cell_value = int(value)
                    except (ValueError, TypeError):
                        cell_value = value

                cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                cell.font = body_font
                cell.border = thin_border
                cell.alignment = left_align if isinstance(cell_value, str) else center_align

        # Auto-adjust column widths
        for col_idx in range(1, len(headers) + 1):
            max_len = len(str(headers[col_idx - 1])) if col_idx <= len(headers) else 0
            for row_idx in range(2, len(rows) + 2):
                val = ws.cell(row=row_idx, column=col_idx).value
                if val is not None:
                    max_len = max(max_len, len(str(val)))
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 50)

        # Freeze header row
        ws.freeze_panes = "A2"

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue(), f"{_safe_stem(title)}.xlsx"


# ── PowerPoint ───────────────────────────────────────────


def build_pptx(
    content: str,
    title: str,
    author_info: str = "",
) -> tuple[bytes, str, str]:
    """
    Build a PowerPoint (.pptx) file with title slide and content slides.

    Parses Markdown content into slides:
    - # Heading → Title slide
    - ## Heading → New slide title
    - ### Heading → Subtitle
    - Bullet points → Bullet list
    - Plain text → Body text
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colors
    PRIMARY = RGBColor(0x44, 0x72, 0xC4)
    DARK = RGBColor(0x1A, 0x1A, 0x1A)
    GRAY = RGBColor(0x55, 0x55, 0x55)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def _add_title_slide(title: str, subtitle: str):
        """Add a title slide."""
        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        # Title text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Subtitle text box
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, line in enumerate(subtitle.strip().split("\n")):
                if i == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                p2.text = line.strip()
                p2.font.size = Pt(18)
                p2.font.color.rgb = GRAY
                p2.alignment = PP_ALIGN.CENTER

    def _add_content_slide(slide_title: str, bullets: list[str]):
        """Add a content slide with title and bullet points."""
        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Bullets
        if bullets:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                # Clean markdown markers
                clean = bullet.strip()
                if clean.startswith("- "):
                    clean = clean[2:]
                elif clean.startswith("  - "):
                    clean = "    " + clean[4:]
                # Bold handling
                p2.text = clean
                p2.font.size = Pt(18)
                p2.font.color.rgb = DARK
                p2.space_after = Pt(8)

    # Parse content into slides
    lines = content.replace("\r\n", "\n").split("\n")
    subtitle = author_info.strip() if author_info else ""

    # First # is title slide
    current_title = title
    current_bullets: list[str] = []
    started = False

    for line in lines:
        stripped = line.strip()

        if stripped.startswith("# ") and not started:
            # Title slide
            current_title = stripped[2:].strip()
            _add_title_slide(current_title, subtitle)
            started = True
            continue

        if not started:
            # Content before first heading → title slide
            _add_title_slide(title, subtitle)
            started = True

        if stripped.startswith("## "):
            # New slide
            if current_bullets:
                _add_content_slide(current_title, current_bullets)
                current_bullets = []
            current_title = stripped[3:].strip()
        elif stripped.startswith("### "):
            current_bullets.append(stripped[4:].strip())
        elif stripped.startswith("- ") or stripped.startswith("* "):
            current_bullets.append(stripped)
        elif stripped and not stripped.startswith("```"):
            current_bullets.append(stripped)

    # Don't forget the last slide
    if current_bullets:
        _add_content_slide(current_title, current_bullets)
    elif not started:
        # No content at all, just title slide
        _add_title_slide(title, subtitle)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue(), f"{_safe_stem(title)}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# ── Main entry ─────────────────────────────────────────────


def build_attachment(
    content: str,
    export_format: str,
    title: str,
    author_info: str = "",
) -> tuple[bytes, str, str]:
    """Build an export attachment in the specified format."""
    fmt = export_format.lower().strip()
    if fmt == "txt":
        return build_txt(content, title, author_info)
    if fmt == "docx":
        return build_docx(content, title, author_info)
    if fmt == "pdf":
        return build_pdf(content, title, author_info)
    if fmt == "pptx":
        return build_pptx(content, title, author_info)
    raise ValueError(f"Unsupported export format: {export_format}")


# Excel 使用独立入口 build_excel（不需要 content/author_info）
# 在 agent_tools.py 中直接调用 build_excel(sheets_data, title)
